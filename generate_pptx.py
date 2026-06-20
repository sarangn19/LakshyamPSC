from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x0F, 0x0F, 0x0F)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
PURPLE_LIGHT = RGBColor(0xC4, 0xB5, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x88)
GRAY_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x34, 0xD3, 0x99)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xFB, 0x71, 0x85)
BORDER = RGBColor(0x2A, 0x2A, 0x3E)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, radius=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(tf, text, font_size=16, color=GRAY_LIGHT, bold=False, space_before=6, alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.space_before = Pt(space_before)
    p.alignment = alignment

def add_card(slide, left, top, width, height, title, body, title_color=PURPLE):
    card = add_rect(slide, left, top, width, height, BG_CARD, BORDER, radius=8)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4), title, 14, title_color, True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.65), body, 12, GRAY_LIGHT)

def slide_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4), f"{num}/{total}", 10, GRAY, alignment=PP_ALIGN.RIGHT)

# ── Helper to add bullet list in a text box ──
def add_bullet_list(slide, left, top, width, height, items, color=GRAY_LIGHT, font_size=12):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_before = Pt(4)
    return txBox

TOTAL = 16

# ═══════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5), "Lakshyam", 72, PURPLE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.8), "Your AI Coach for Kerala PSC Exams", 28, RGBColor(0xA7,0x8B,0xFA), False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5), "AI-native exam preparation platform  ·  iOS  ·  Android  ·  Web", 16, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.4), "Investor Pitch  |  June 2026", 12, RGBColor(0x55,0x55,0x55), False, PP_ALIGN.CENTER)
slide_number(slide, 1, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 2: Problem
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "The Problem", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), "1M+ aspirants, zero modern tools", 36, WHITE, True)
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(2.5), "📊  Massive Market", "1M+ applicants every year for Kerala PSC exams. Lakhs compete for hundreds of vacancies. No digital platform exists for this segment.")
add_card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.5), "📚  Outdated Solutions", "Printed books, coaching centers, and generic exam apps. Nothing built specifically for Kerala PSC syllabus or Malayalam-speaking users.")
add_card(slide, Inches(0.8), Inches(4.6), Inches(5.8), Inches(2.5), "🧠  One-Size-Fits-None", "Static question banks. No personalization. No adaptive learning. Every student, regardless of ability, gets the same content.")
add_card(slide, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.5), "🌐  Language & Access Gap", "Most tools are English-only. Kerala aspirants need Malayalam support. Rural users need offline access. Existing solutions fail on all fronts.")
slide_number(slide, 2, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 3: Solution
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "The Solution", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), "Full-stack AI learning platform", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4), "Fresh AI-generated questions  ·  Cognitive twin  ·  Adaptive engine  ·  Spaced repetition  ·  Offline-first", 14, RGBColor(0x88,0x88,0x88))

cards = [
    ("🤖  AI Question Generator", "Fresh questions for any subject, topic, or difficulty — in English or Malayalam. Multi-model pipeline (Groq \u2192 Gemini \u2192 OpenRouter)."),
    ("🧬  Cognitive Twin", "A digital replica of each learner's knowledge. Tracks mastery, forgetting, confidence, and gaps across 47+ subtopics with prerequisite chains."),
    ("🎯  Adaptive Engine", "Bayesian Knowledge Tracing selects optimal next question. Real-time difficulty adjustment. Learns from every answer."),
    ("⏰  Spaced Repetition", "SM-2 algorithm schedules reviews at optimal intervals (1, 3, 7, 14, 30, 60, 120, 240 days). Retention dashboard."),
    ("💬  AI Tutor", "Kerala PSC context-aware chatbot. Understands Malayalam. Explains concepts, generates practice questions on demand."),
    ("📰  Current Affairs", "Auto-fetched daily from NewsAPI. Categorized (Kerala / National / Schemes / Appointments / Awards). With images."),
]
for i, (title, body) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(2.1 + row * 2.6)
    add_card(slide, x, y, Inches(3.8), Inches(2.3), title, body)
slide_number(slide, 3, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 4: How It Works
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "How It Works", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), "Every answer trains the model", 36, WHITE, True)

# Flow diagram boxes
flow_items = ["Student answers\nquestion", "Cognitive Twin\nupdates", "Adaptive Engine\npicks next topic", "AI generates\nfresh question", "Validated\n& served", "Spaced Rep\nschedules review"]
box_w = Inches(1.8)
box_h = Inches(1.0)
gap = Inches(0.15)
start_x = Inches(0.8)
start_y = Inches(2.5)
for i, item in enumerate(flow_items):
    x = start_x + i * (box_w + gap)
    shape = add_rect(slide, x, start_y, box_w, box_h, BG_CARD, PURPLE, radius=8)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = PURPLE_LIGHT
    p.font.name = 'Calibri'
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    if i < len(flow_items) - 1:
        add_text_box(slide, x + box_w, start_y + Inches(0.3), Inches(0.2), Inches(0.4), "\u2192", 18, GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.8), "\"The system gets smarter with every answer. No two students see the same path.\"", 20, PURPLE_LIGHT, False, PP_ALIGN.CENTER)
slide_number(slide, 4, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 5: Cognitive Twin
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "Core Differentiator", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), "The Cognitive Twin", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4), "A persistent knowledge model across a 3-level syllabus with prerequisite chains", 16, GRAY)

add_card(slide, Inches(0.8), Inches(2.1), Inches(5.8), Inches(2.2), "🧠  Syllabus Structure", "9 Subjects \u2192 20+ Topics \u2192 47+ Subtopics\nEach node has prerequisites, mastery state, and prerequisite-chain validation.")
add_card(slide, Inches(6.8), Inches(2.1), Inches(5.8), Inches(2.2), "📊  Per-Node Tracking", "Mastery score, accuracy, hesitation, forgetting rate\nTrend direction (improving / declining / stable)\n7-day, 30-day, 90-day retention rates")
add_card(slide, Inches(0.8), Inches(4.6), Inches(5.8), Inches(2.5), "🔄  Gap Lifecycle", "Open \u2192 Improving \u2192 Closing \u2192 Closed \u2192 Retained\nEach gap tracked with lifecycle status, confidence calibration, and recommended intervention count.")
add_card(slide, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.5), "🎯  What This Enables", "Know exactly what each student is about to forget\nPrioritize gaps that matter for their target exam\nMeasure real learning progress, not just question count\nWeighted recommendations combining weakness, forgetting, exam blueprint, and learner stage.")
slide_number(slide, 5, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 6: AI Pipeline
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "Infrastructure", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), "AI Multi-Model Pipeline", 36, WHITE, True)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4), "No single point of failure. No expensive vendor lock-in. $0 AI inference cost at current scale.", 14, GREEN)

# Pipeline visual
providers = [
    ("Groq — Primary", "Llama 3.1 8B  |  30 req/min free", PURPLE),
    ("OpenRouter — Fallback 1", "Gemini 2.0 Flash Lite", YELLOW),
    ("Gemini Direct — Fallback 2", "Gemini 2.0 Flash  |  1,500 req/day free", GREEN),
]
for i, (name, desc, color) in enumerate(providers):
    y = Inches(2.2 + i * 1.3)
    shape = add_rect(slide, Inches(3.5), y, Inches(6.5), Inches(1.0), BG_CARD, color, radius=8)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(18)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    add_paragraph(tf, desc, 13, GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
    if i < len(providers) - 1:
        add_text_box(slide, Inches(5.8), y + Inches(1.0), Inches(2), Inches(0.4), "\u25bc  fallback", 11, GRAY, False, PP_ALIGN.CENTER)

# Stats row
stats = [("1,500+", "Questions/day free"), ("3", "AI providers"), ("$0", "Inference cost")]
for i, (num, label) in enumerate(stats):
    x = Inches(1.5 + i * 3.8)
    add_text_box(slide, x, Inches(5.5), Inches(3), Inches(0.6), num, 36, PURPLE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.1), Inches(3), Inches(0.4), label, 12, GRAY, False, PP_ALIGN.CENTER)
slide_number(slide, 6, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 7: Screens
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "Product", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), "Student Screens", 36, WHITE, True)

screens = [
    ("📊  Dashboard", "Progress, accuracy, streak, due reviews, subject breakdown"),
    ("🎯  MCQ Engine", "AI-generated questions, confidence tracking, timer, flag/report"),
    ("🃏  Flashcards", "SM-2 spaced repetition, flip & rate, mastery tracking"),
    ("🗺️  Knowledge Map", "3-level tree with mastery heatmap, prerequisite warnings"),
    ("💬  AI Tutor", "Malayalam/English chatbot, explains concepts, generates practice"),
    ("📈  Analytics", "Gap lifecycle, retention curves, confidence calibration"),
    ("📰  Current Affairs", "Daily news with images, categorized, exam-relevant"),
    ("🏆  Retention Dashboard", "Heat maps, sparklines, due/at-risk subtopics"),
    ("🎯  Goal Tracker", "Exam roadmaps with phases, daily targets, study streak"),
]
for i, (title, body) in enumerate(screens):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.7 + row * 1.9)
    add_card(slide, x, y, Inches(3.8), Inches(1.6), title, body)
slide_number(slide, 7, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 8: Admin Portals
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "Platform", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), "Admin & Superadmin Portals", 36, WHITE, True)

# Admin portal
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(3.5), "🛠️  Admin Portal — 8 Tabs", "")
admin_items = [
    "Dashboard — Platform metrics and overview",
    "Question Management — Full CRUD for question bank",
    "Current Affairs Publishing — Schedule and publish CA",
    "Content Quality Center — Review pipeline",
    "Question Audit — Flagged question review",
    "Bulk Upload — Batch question import",
    "Learner Support — Ticket management system",
    "Learning Analytics — Aggregate learner analytics",
]
add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.3), Inches(3.0), admin_items, GRAY_LIGHT, 11)

# Superadmin portal
add_card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(3.5), "⚡  Superadmin Portal — 8 Tabs", "")
superadmin_items = [
    "Executive Dashboard — Top-level system KPIs",
    "Cognitive Twin Config — Tune weights (weakness, forgetting, confusion, coverage)",
    "Recommendation Analytics — Monitor recommendation quality",
    "User & Role Management — User CRUD, role assignment",
    "Access Control — Role & permission management",
    "System Monitoring — DB health, sync failures, storage",
    "Experiment Center — A/B testing framework",
    "Audit Logs — Full admin action audit trail",
]
add_bullet_list(slide, Inches(7.1), Inches(2.3), Inches(5.3), Inches(3.0), superadmin_items, GRAY_LIGHT, 11)

add_text_box(slide, Inches(3.5), Inches(5.8), Inches(6), Inches(0.5), "Student \u2192  Admin \u2192  Superadmin", 16, PURPLE_LIGHT, False, PP_ALIGN.CENTER)
slide_number(slide, 8, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 9: Technology
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "Architecture", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), "Technology Stack", 36, WHITE, True)

tech_cards = [
    ("🎨  Frontend", "React Native 0.76 + Expo 52  ·  TypeScript\nZustand with AsyncStorage offline persistence\nCustom i18n (English + Malayalam)"),
    ("⚙️  Backend", "Supabase (PostgreSQL, Auth, Edge Functions)\n16 database migrations, 8 Deno edge functions\nRealtime subscriptions for sync"),
    ("🧠  AI & ML", "Groq + Gemini + OpenRouter + OpenAI\nBayesian Knowledge Tracing (BKT)\nCognitive twin with gap lifecycle tracking\nSM-2 spaced repetition algorithm"),
    ("🌐  Platform & Scale", "iOS + Android + Web — single codebase\nVercel (Web), Expo (Mobile)\n~110 source files, 30+ services, 12 stores\nOffline-first with background sync"),
]
for i, (title, body) in enumerate(tech_cards):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.3)
    y = Inches(1.7 + row * 2.8)
    add_card(slide, x, y, Inches(5.8), Inches(2.5), title, body)
slide_number(slide, 9, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 10: Market
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "Market", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), "Massive underserved opportunity", 36, WHITE, True)

stats_data = [
    ("1M+", "Annual PSC applicants"),
    ("5+", "Exam types"),
    ("10,000+", "Coaching centers in Kerala"),
    ("96%", "Kerala literacy rate"),
    ("80%+", "Smartphone penetration"),
    ("0", "Dedicated Kerala PSC apps"),
]
for i, (num, label) in enumerate(stats_data):
    col = i % 3
    row = i // 3
    x = Inches(1.5 + col * 3.8)
    y = Inches(2.0 + row * 2.5)
    shape = add_rect(slide, x, y, Inches(3.2), Inches(2.0), BG_CARD, BORDER, radius=8)
    add_text_box(slide, x, y + Inches(0.3), Inches(3.2), Inches(0.8), num, 42, PURPLE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.1), Inches(3.2), Inches(0.6), label, 13, GRAY_LIGHT, False, PP_ALIGN.CENTER)
slide_number(slide, 10, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 11: Competition
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "Competition", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), "No one else does this", 36, WHITE, True)

# Table
rows_data = [
    ("AI-generated questions", "✅", "❌", "❌", "❌"),
    ("Adaptive learning", "✅", "❌", "❌", "❌"),
    ("Cognitive twin", "✅", "❌", "❌", "❌"),
    ("Malayalam support", "✅", "❌", "✅", "✅"),
    ("Offline-first", "✅", "❌", "❌", "✅"),
    ("Multi-platform", "✅", "❌", "❌", "❌"),
    ("Real-time analytics", "✅", "❌", "❌", "❌"),
    ("Cost", "Free tier", "Paid", "₹5K-50K", "₹500-2K"),
]
headers = ["Capability", "Lakshyam", "Generic Apps", "Coaching", "Books"]
col_widths = [Inches(3.5), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
x_start = Inches(1.0)
y_start = Inches(1.8)
row_h = Inches(0.55)

# Header row
x = x_start
for j, (h, w) in enumerate(zip(headers, col_widths)):
    shape = add_rect(slide, x, y_start, w, row_h, PURPLE, None)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = h
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    x += w

# Data rows
for i, row_data in enumerate(rows_data):
    x = x_start
    y = y_start + row_h * (i + 1)
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
    for j, (val, w) in enumerate(zip(row_data, col_widths)):
        shape = add_rect(slide, x, y, w, row_h, bg, BORDER)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(12)
        p.font.color.rgb = PURPLE if j == 1 and val == "✅" else (GREEN if j == 1 and "Free" in val else GRAY_LIGHT)
        p.font.bold = (j == 1)
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
        x += w
slide_number(slide, 11, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 12: Business Model
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "Monetization", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), "Business Model", 36, WHITE, True)

biz_cards = [
    ("🎁  Free", "AI-generated questions\nBasic analytics\nDaily current affairs", "₹0", GREEN),
    ("⭐  Premium", "Unlimited questions\nOffline question bank\nPriority AI generation\nDetailed analytics\nAd-free experience", "₹199/mo", YELLOW),
    ("🚀  Pro", "All Premium features\nUnlimited AI tutor\nCognitive twin deep dive\nStudy groups & leaderboards\nPriority support", "₹499/mo", RED),
]
for i, (title, features, price, color) in enumerate(biz_cards):
    x = Inches(0.8 + i * 4.1)
    add_card(slide, x, Inches(1.8), Inches(3.8), Inches(3.5), title, features)
    add_text_box(slide, x, Inches(5.3), Inches(3.8), Inches(0.5), price, 24, color, True, PP_ALIGN.CENTER)

add_card(slide, Inches(0.8), Inches(5.8), Inches(12), Inches(1.2), "🏢  Enterprise (planned)", "Coaching center white-label solution  ·  Batch management  ·  Custom content  ·  Analytics dashboard  ·  Custom pricing\n\nAI cost per question: ~₹0.001 (fraction of a paisa)  |  Coaching centers charge ₹15,000–50,000/year per student")
slide_number(slide, 12, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 13: Revenue Model
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "Revenue Model", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.5), "How Lakshyam makes money", 32, WHITE, True)

# Unit economics
add_card(slide, Inches(0.8), Inches(1.4), Inches(5.8), Inches(1.5), "📊  Unit Economics", "")
unit_items = [
    "AI inference cost per question: ₹0.001 (free tier Groq+Gemini)",
    "Monthly active user ~300 questions = ₹0.30 AI cost (free → ₹0)",
    "Premium ₹199/mo → 99.85% gross margin with free APIs",
    "Upgrade to GPT-4o-mini adds ₹5/user/mo → 97.5% margin still",
    "Upgrade to GPT-4o/Claude adds ₹80/user/mo → profitable at ₹299+",
]
add_bullet_list(slide, Inches(1.1), Inches(1.8), Inches(5.3), Inches(1.3), unit_items, GRAY_LIGHT, 11)

add_card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(1.5), "💰  Revenue Streams", "")
rev_items = [
    "B2C Subscriptions: Premium (₹199/mo) + Pro (₹499/mo)",
    "B2B Enterprise: Coaching center white-label licensing",
    "Content Marketplace: Curated question sets & mock tests",
    "Future: Ad revenue from free tier, affiliate (books/courses)",
]
add_bullet_list(slide, Inches(7.1), Inches(1.8), Inches(5.3), Inches(1.3), rev_items, GRAY_LIGHT, 11)

# Projections table
add_rect(slide, Inches(0.8), Inches(3.1), Inches(11.8), Inches(0.45), PURPLE, None)
add_text_box(slide, Inches(1.0), Inches(3.15), Inches(6), Inches(0.4), "Projected Revenue at Scale", 13, WHITE, True)

proj_headers = ["Metric", "Year 1", "Year 2", "Year 3"]
proj_col_w = [Inches(4.0), Inches(2.6), Inches(2.6), Inches(2.6)]
proj_x = Inches(0.8)
proj_y = Inches(3.55)
proj_row_h = Inches(0.33)

for j, (h, w) in enumerate(zip(proj_headers, proj_col_w)):
    x = proj_x + sum(proj_col_w[k] for k in range(j))
    shape = add_rect(slide, x, proj_y, w, proj_row_h, BG_CARD, BORDER)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = h
    p.font.size = Pt(9)
    p.font.color.rgb = PURPLE_LIGHT
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

proj_data = [
    ("Free users", "10,000", "50,000", "200,000"),
    ("Conversion rate", "3%", "5%", "8%"),
    ("Paying users", "300", "2,500", "16,000"),
    ("Avg. revenue/user/mo", "₹199", "₹249", "₹299"),
    ("Monthly recurring revenue", "₹60K", "₹6.2L", "₹47.8L"),
    ("Annual run rate (ARR)", "₹7.2L", "₹74.4L", "₹5.7Cr"),
]
for i, row in enumerate(proj_data):
    y = proj_y + proj_row_h * (i + 1)
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x15, 0x15, 0x28)
    for j, (val, w) in enumerate(zip(row, proj_col_w)):
        x = proj_x + sum(proj_col_w[k] for k in range(j))
        shape = add_rect(slide, x, y, w, proj_row_h, bg, BORDER)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(9)
        p.font.color.rgb = GREEN if j > 0 and ("L" in val or "Cr" in val) else (YELLOW if j == 0 else GRAY_LIGHT)
        p.font.bold = True if j == 0 else False
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER

# API cost rows (subdued)
api_label_y = proj_y + proj_row_h * (len(proj_data) + 1)
api_data = [
    ("API cost (free tier)", "₹0", "₹0", "₹0"),
    ("API cost (GPT-4o-mini)", "₹1.5K", "₹12.5K", "₹80K"),
    ("API cost (GPT-4o/Claude)", "₹24K", "₹2L", "₹12.8L"),
]
for i, row in enumerate(api_data):
    y = api_label_y + proj_row_h * i
    bg = RGBColor(0x10, 0x10, 0x20)
    for j, (val, w) in enumerate(zip(row, proj_col_w)):
        x = proj_x + sum(proj_col_w[k] for k in range(j))
        shape = add_rect(slide, x, y, w, proj_row_h, bg, BORDER)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(8)
        p.font.italic = True
        p.font.color.rgb = GRAY if "free" in val.lower() else (GREEN if "free" in row[0].lower() else (YELLOW if "mini" in row[0].lower() else RGBColor(0xFB, 0x71, 0x85)))
        p.font.bold = False
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER

# Summary line
add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.5), "Break-even: ~500 paid users (₹199/mo = ₹99.5K/mo). All API tiers remain highly profitable. Free tier sustains launch at zero inference cost.", 10, GRAY)
slide_number(slide, 13, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 14: Traction
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "Progress", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), "Traction & Milestones", 36, WHITE, True)

milestones = [
    ("✅  AI Question Generation", "Multi-model pipeline working. 1,500+ questions/day free capacity."),
    ("✅  Cognitive Twin", "Full implementation with gap lifecycle, retention tracking, exam weights."),
    ("✅  Adaptive Engine", "BKT-based with real-time difficulty, prerequisite chains, learner stages."),
    ("✅  Admin Portal", "Complete content management, quality center, learner support, analytics."),
    ("✅  Current Affairs", "Daily auto-fetch from NewsAPI. Categorized, with images. Live on Vercel."),
    ("✅  Spaced Repetition", "SM-2 algorithm. Retention dashboard with heat maps and sparklines."),
]
for i, (title, body) in enumerate(milestones):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.7 + row * 2.0)
    add_card(slide, x, y, Inches(3.8), Inches(1.7), title, body)

# Pipeline
add_text_box(slide, Inches(0.8), Inches(5.9), Inches(2), Inches(0.4), "Next:", 14, PURPLE, True)
add_rect(slide, Inches(2.5), Inches(5.9), Inches(10), Inches(0.6), BG_CARD, YELLOW, 6)
add_text_box(slide, Inches(2.7), Inches(5.95), Inches(9.5), Inches(0.5), "Play Store / App Store launch  ·  Premium features  ·  Coaching center partnerships  ·  Other state PSC exams", 13, YELLOW, False)
slide_number(slide, 14, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 15: The Ask
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4), "The Ask", 14, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6), "Building the definitive AI learning platform for India's exam market", 32, WHITE, True)

add_card(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(2.5), "💰  Seed Investment", "")
add_text_box(slide, Inches(1.2), Inches(2.3), Inches(5), Inches(1.0), "$50,000 – $100,000", 42, PURPLE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.2), Inches(3.1), Inches(5), Inches(1.0), "Marketing & user acquisition in Kerala\nContent partnerships with coaching centers\nPremium feature development\nPlay Store / App Store launch", 13, GRAY_LIGHT)

add_card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.5), "🎯  What You Get", "")
investor_items = [
    "Complete product: ~110 source files, full AI pipeline, admin backend",
    "Modern stack: React Native + Supabase + Vercel — low burn rate",
    "First-mover advantage in a 1M+ user market with zero competitors",
    "Data moat: more users = smarter AI = stronger lock-in",
    "Solo founder with full-stack + AI expertise",
]
add_bullet_list(slide, Inches(7.1), Inches(2.3), Inches(5.3), Inches(2.0), investor_items, GRAY_LIGHT, 12)

add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(1.0), '"Lakshyam" means "target" in Malayalam.\nWe help every aspirant hit theirs.', 24, PURPLE_LIGHT, False, PP_ALIGN.CENTER)
slide_number(slide, 15, TOTAL)

# ═══════════════════════════════════════════
# SLIDE 16: Thank You
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5), "Thank You", 72, PURPLE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.8), "Lakshyam — Your AI Coach for Kerala PSC Exams", 28, PURPLE_LIGHT, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.5), "Built by a solo developer  ·  React Native + Supabase + Vercel", 16, GRAY, False, PP_ALIGN.CENTER)
slide_number(slide, 16, TOTAL)

# Save
prs.save('Lakshyam_Investor_Pitch.pptx')
print("Saved: Lakshyam_Investor_Pitch.pptx")
